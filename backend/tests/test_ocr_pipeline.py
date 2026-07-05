"""Unit tests for the extraction dispatch pipeline (no ML deps required)."""

from io import BytesIO
from unittest.mock import AsyncMock, patch

import pytest

from app.workers.ocr.pipeline import (
    DOCX_MIME,
    PPTX_MIME,
    decode_text_bytes,
    extract_text_from_bytes,
    extract_text_from_docx,
    extract_text_from_pptx,
)


class TestDecodeTextBytes:
    def test_decodes_utf8(self):
        assert decode_text_bytes(b"hello") == "hello"

    def test_handles_invalid_utf8_via_replacement(self):
        # 0xff is not valid in UTF-8 — must not raise.
        out = decode_text_bytes(b"hello \xff world")
        assert "hello" in out and "world" in out


class TestExtractTextFromBytes:
    async def test_plain_text(self):
        assert await extract_text_from_bytes(b"hello world", "text/plain") == "hello world"

    async def test_markdown_text(self):
        assert await extract_text_from_bytes(b"# hi", "text/markdown; charset=utf-8") == "# hi"

    async def test_audio_dispatches_to_transcription(self):
        with patch(
            "app.transcription.whisper.transcribe_audio_bytes",
            return_value="hello from whisper",
        ) as transcribe:
            out = await extract_text_from_bytes(b"\x00\x01\x02\x03", "audio/mpeg")
            transcribe.assert_called_once()
            assert out == "hello from whisper"

    async def test_unknown_mime_raises(self):
        with pytest.raises(ValueError):
            await extract_text_from_bytes(b"", "application/x-msdownload")

    async def test_image_combines_ocr_and_vision(self):
        with (
            patch(
                "app.workers.ocr.engines.ocr_image_bytes",
                return_value="OCR-TEXT-HERE",
            ) as ocr,
            patch(
                "app.vision.describe.describe_image_bytes",
                new=AsyncMock(return_value="VISION-DESC-HERE"),
            ),
        ):
            out = await extract_text_from_bytes(b"\xff\xd8\xff\xe0", "image/jpeg")

        ocr.assert_called_once()
        assert "--- OCR text ---" in out
        assert "OCR-TEXT-HERE" in out
        assert "--- Visual description ---" in out
        assert "VISION-DESC-HERE" in out

    async def test_image_continues_when_vision_fails(self):
        from app.vision.describe import VisionError

        with (
            patch(
                "app.workers.ocr.engines.ocr_image_bytes",
                return_value="just-the-ocr",
            ),
            patch(
                "app.vision.describe.describe_image_bytes",
                new=AsyncMock(side_effect=VisionError("upstream down")),
            ),
        ):
            out = await extract_text_from_bytes(b"\xff", "image/png")

        assert "just-the-ocr" in out
        assert "Visual description" not in out

    async def test_image_continues_when_ocr_empty(self):
        with (
            patch(
                "app.workers.ocr.engines.ocr_image_bytes",
                return_value="",
            ),
            patch(
                "app.vision.describe.describe_image_bytes",
                new=AsyncMock(return_value="only the vision description"),
            ),
        ):
            out = await extract_text_from_bytes(b"\xff", "image/png")

        assert "only the vision description" in out
        assert "OCR text" not in out

    async def test_pdf_dispatches_to_pdf_path(self):
        with patch(
            "app.workers.ocr.pipeline.extract_text_from_pdf",
            return_value="PDF-MOCK",
        ) as fn:
            out = await extract_text_from_bytes(b"%PDF-1.4", "application/pdf")
            fn.assert_called_once()
            assert out == "PDF-MOCK"

    async def test_docx_dispatches_to_docx_path(self):
        with patch(
            "app.workers.ocr.pipeline.extract_text_from_docx",
            return_value="DOCX-MOCK",
        ) as fn:
            out = await extract_text_from_bytes(b"PK\x03\x04", DOCX_MIME)
            fn.assert_called_once()
            assert out == "DOCX-MOCK"

    async def test_pptx_dispatches_to_pptx_path(self):
        with patch(
            "app.workers.ocr.pipeline.extract_text_from_pptx",
            return_value="PPTX-MOCK",
        ) as fn:
            out = await extract_text_from_bytes(b"PK\x03\x04", PPTX_MIME)
            fn.assert_called_once()
            assert out == "PPTX-MOCK"


def _make_tiny_docx() -> bytes:
    """Build a minimal in-memory .docx with a paragraph, a heading, and a
    2x2 table so we can assert real python-docx extraction end-to-end.
    Skips the calling test when python-docx isn't installed — the module
    lives in the [worker] extras, which CI's unit-test job doesn't pull."""
    pytest.importorskip("docx")
    from docx import Document  # type: ignore[import-not-found]

    doc = Document()
    doc.add_heading("Report title", level=1)
    doc.add_paragraph("First paragraph body.")
    doc.add_paragraph("")  # blank — must be filtered
    doc.add_paragraph("Second paragraph body.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Header A"
    table.rows[0].cells[1].text = "Header B"
    table.rows[1].cells[0].text = "Row1 A"
    table.rows[1].cells[1].text = "Row1 B"
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_tiny_pptx(*, include_notes: bool = True) -> bytes:
    """Skips the calling test when python-pptx isn't installed — same
    reason as _make_tiny_docx."""
    pytest.importorskip("pptx")
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for idx in range(2):
        slide = prs.slides.add_slide(blank_layout)
        tb = slide.shapes.add_textbox(0, 0, 500, 500)
        tb.text_frame.text = f"Slide {idx + 1} body content"
        if include_notes:
            slide.notes_slide.notes_text_frame.text = f"Notes for slide {idx + 1}"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestDocxExtraction:
    def test_extracts_paragraphs_headings_and_tables(self):
        out = extract_text_from_docx(_make_tiny_docx())
        assert "Report title" in out
        assert "First paragraph body." in out
        assert "Second paragraph body." in out
        # Table rows joined with " | " separator.
        assert "Header A | Header B" in out
        assert "Row1 A | Row1 B" in out
        # Blank paragraph filtered out — no empty double-newlines beyond spec.
        assert "\n\n\n" not in out

    async def test_dispatch_returns_the_same_text_as_direct_extraction(self):
        # No mock — the dispatch actually invokes the real docx path.
        raw = _make_tiny_docx()
        direct = extract_text_from_docx(raw)
        via_dispatch = await extract_text_from_bytes(raw, DOCX_MIME)
        assert direct == via_dispatch


class TestPptxExtraction:
    def test_extracts_body_and_speaker_notes_per_slide(self):
        out = extract_text_from_pptx(_make_tiny_pptx(include_notes=True))
        assert "--- slide 1 ---" in out
        assert "--- slide 2 ---" in out
        assert "Slide 1 body content" in out
        assert "Slide 2 body content" in out
        assert "[speaker notes]" in out
        assert "Notes for slide 1" in out
        assert "Notes for slide 2" in out

    def test_slides_without_notes_have_no_notes_section(self):
        out = extract_text_from_pptx(_make_tiny_pptx(include_notes=False))
        assert "Slide 1 body content" in out
        assert "[speaker notes]" not in out
